"""PPT Reader - Extract text content from PowerPoint files"""

from dataclasses import dataclass, field
from pathlib import Path
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.table import Table
from pptx.util import Pt
import logging

logger = logging.getLogger(__name__)


@dataclass
class TextRun:
    """A run of text with consistent formatting"""
    text: str
    font_name: str | None = None
    font_size: Pt | None = None
    bold: bool | None = None
    italic: bool | None = None


@dataclass
class TextElement:
    """A text element from a shape"""
    slide_index: int
    shape_id: int
    paragraph_index: int
    run_index: int
    text: str
    shape_type: str = "textbox"


@dataclass
class SlideContent:
    """Content extracted from a single slide"""
    index: int
    title: str | None = None
    elements: list[TextElement] = field(default_factory=list)


@dataclass
class PresentationContent:
    """Content extracted from a presentation"""
    file_path: Path
    slides: list[SlideContent] = field(default_factory=list)

    @property
    def all_texts(self) -> list[str]:
        """Get all unique texts for translation"""
        texts = set()
        for slide in self.slides:
            for element in slide.elements:
                if element.text and element.text.strip():
                    texts.add(element.text)
        return list(texts)


class PPTReader:
    """Read and extract text from PowerPoint files"""

    def __init__(self, file_path: str | Path):
        self.file_path = Path(file_path)
        if not self.file_path.exists():
            raise FileNotFoundError(f"File not found: {self.file_path}")
        if not self.file_path.suffix.lower() == ".pptx":
            raise ValueError("Only .pptx files are supported")

        self.presentation = Presentation(str(self.file_path))

    def extract_content(self) -> PresentationContent:
        """Extract all text content from the presentation"""
        content = PresentationContent(file_path=self.file_path)

        for slide_idx, slide in enumerate(self.presentation.slides):
            slide_content = SlideContent(index=slide_idx)

            # Try to get slide title
            if slide.shapes.title:
                slide_content.title = slide.shapes.title.text

            # Extract text from all shapes
            for shape in slide.shapes:
                self._extract_shape_text(shape, slide_idx, slide_content)

            content.slides.append(slide_content)

        return content

    def _extract_shape_text(
        self, shape: BaseShape, slide_idx: int, slide_content: SlideContent
    ) -> None:
        """Extract text from a shape recursively"""
        # Handle grouped shapes
        if isinstance(shape, GroupShape):
            for child_shape in shape.shapes:
                self._extract_shape_text(child_shape, slide_idx, slide_content)
            return

        # Handle tables
        if shape.has_table:
            self._extract_table_text(shape.table, slide_idx, shape.shape_id, slide_content)
            return

        # Handle text frames - extract entire text frame as one element
        if shape.has_text_frame:
            # Replace vertical tab (soft line break) with newline
            full_text = shape.text_frame.text.replace('\x0b', '\n')
            if full_text and full_text.strip():
                element = TextElement(
                    slide_index=slide_idx,
                    shape_id=shape.shape_id,
                    paragraph_index=0,
                    run_index=0,
                    text=full_text,
                    shape_type="textbox",
                )
                slide_content.elements.append(element)

    def _extract_table_text(
        self, table: Table, slide_idx: int, shape_id: int, slide_content: SlideContent
    ) -> None:
        """Extract text from a table"""
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if cell.text and cell.text.strip():
                    # Use a composite index for table cells
                    element = TextElement(
                        slide_index=slide_idx,
                        shape_id=shape_id,
                        paragraph_index=row_idx * 1000 + col_idx,  # Encode row/col
                        run_index=0,
                        text=cell.text,
                        shape_type="table",
                    )
                    slide_content.elements.append(element)

    def get_presentation(self) -> Presentation:
        """Get the underlying Presentation object"""
        return self.presentation
