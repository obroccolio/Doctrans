"""PPT Writer - Write translated content back to PowerPoint"""

import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.group import GroupShape
from pptx.dml.color import RGBColor
from pptx.util import Pt, Emu
import logging

from .translator import is_math_content

logger = logging.getLogger(__name__)


class PPTWriter:
    """Write translated content back to PowerPoint files"""

    def __init__(
        self,
        presentation: Presentation,
        chinese_font: str = "Microsoft YaHei",
        llm_client=None,
        render_math: bool = True,
    ):
        self.presentation = presentation
        self.chinese_font = chinese_font
        self.llm_client = llm_client
        self.render_math = render_math and llm_client is not None
        self._temp_dir = None
        self._math_images: dict[str, Path] = {}

    def apply_translations(self, translations: dict[str, str]) -> None:
        """Apply translations to the presentation"""
        # Pre-render all math formulas
        if self.render_math:
            self._prerender_math_formulas(translations)

        for slide in self.presentation.slides:
            shapes_to_process = list(slide.shapes)
            for shape in shapes_to_process:
                self._apply_to_shape(slide, shape, translations)

    def _prerender_math_formulas(self, translations: dict[str, str]) -> None:
        """Pre-render all math formulas to images."""
        from .math_renderer import render_math_text

        # Create temp directory for math images
        self._temp_dir = tempfile.mkdtemp(prefix="doctrans_math_")
        temp_path = Path(self._temp_dir)

        for text, translated in translations.items():
            # If unchanged and is math content, render it
            if text == translated and is_math_content(text):
                image_path = render_math_text(text, self.llm_client, temp_path)
                if image_path:
                    self._math_images[text] = image_path
                    logger.info(f"Pre-rendered math: '{text}' -> {image_path}")

    def _apply_to_shape(self, slide, shape: BaseShape, translations: dict[str, str]) -> None:
        """Apply translations to a shape recursively"""
        # Handle grouped shapes
        if isinstance(shape, GroupShape):
            for child_shape in shape.shapes:
                self._apply_to_shape(slide, child_shape, translations)
            return

        # Handle tables
        if shape.has_table:
            self._apply_to_table(shape.table, translations)
            return

        # Handle text frames - match by full text frame content
        if shape.has_text_frame:
            full_text = shape.text_frame.text
            if full_text in translations:
                translated_text = translations[full_text]

                # Check if this is a math formula that should be replaced with image
                if full_text in self._math_images:
                    self._replace_shape_with_image(slide, shape, self._math_images[full_text])
                    return

                # Skip if no change (math formulas without rendering) - preserve original
                if translated_text == full_text:
                    return

                # Split translated text by original paragraph structure
                original_paragraphs = [p.text for p in shape.text_frame.paragraphs]
                translated_lines = translated_text.split('\n')

                # Apply translation preserving paragraph structure
                line_idx = 0
                for paragraph in shape.text_frame.paragraphs:
                    if line_idx < len(translated_lines):
                        # Collect all runs text and replace with translated line
                        if paragraph.runs:
                            # Put all translated text in first run, clear others
                            paragraph.runs[0].text = translated_lines[line_idx]
                            if self._contains_chinese(translated_lines[line_idx]):
                                paragraph.runs[0].font.name = self.chinese_font
                            for run in paragraph.runs[1:]:
                                run.text = ""
                        line_idx += 1

    def _replace_shape_with_image(self, slide, shape: BaseShape, image_path: Path) -> None:
        """Replace a text shape with a rendered math image."""
        try:
            # Get the position and size of the original shape
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            # Add the image at the same position
            from PIL import Image
            with Image.open(image_path) as img:
                img_width, img_height = img.size

            # Calculate scaling to fit within original bounds while maintaining aspect ratio
            aspect_ratio = img_width / img_height
            original_aspect = width / height if height > 0 else 1

            if aspect_ratio > original_aspect:
                # Image is wider, fit to width
                new_width = width
                new_height = int(width / aspect_ratio)
            else:
                # Image is taller, fit to height
                new_height = height
                new_width = int(height * aspect_ratio)

            # Center the image in the original position
            new_left = left + (width - new_width) // 2
            new_top = top + (height - new_height) // 2

            # Add the picture
            slide.shapes.add_picture(
                str(image_path),
                new_left,
                new_top,
                new_width,
                new_height
            )

            # Hide the original shape by clearing its text
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = ""

            logger.info(f"Replaced math shape with image: {image_path}")
        except Exception as e:
            logger.error(f"Failed to replace shape with image: {e}")
            # Keep original shape if replacement fails

    def _apply_to_table(self, table, translations: dict[str, str]) -> None:
        """Apply translations to a table"""
        for row in table.rows:
            for cell in row.cells:
                if cell.text in translations:
                    # For tables, we need to preserve cell formatting
                    original_text = cell.text
                    translated_text = translations[original_text]

                    # Clear and rewrite with translated text
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                # Find the matching translation for this run
                                if run.text in translations:
                                    run.text = translations[run.text]
                                    if self._contains_chinese(run.text):
                                        run.font.name = self.chinese_font

    def _contains_chinese(self, text: str) -> bool:
        """Check if text contains Chinese characters"""
        for char in text:
            if '\u4e00' <= char <= '\u9fff':
                return True
        return False

    def save(self, output_path: str | Path) -> Path:
        """Save the presentation to a file"""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(str(output_path))
        logger.info(f"Saved translated presentation to: {output_path}")
        return output_path


def create_output_path(input_path: Path, suffix: str = "_translated") -> Path:
    """Create output path based on input path"""
    return input_path.parent / f"{input_path.stem}{suffix}{input_path.suffix}"
